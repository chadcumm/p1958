#!/usr/bin/env python3
"""Generate MHA PDS HSE Bundle Issue presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY = RGBColor(0x99, 0x99, 0x99)
RED = RGBColor(0xE7, 0x4C, 0x3C)
RED_LIGHT = RGBColor(0xFF, 0x6B, 0x6B)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
GREEN_DARK = RGBColor(0x27, 0xAE, 0x60)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
BLUE_LIGHT = RGBColor(0x60, 0xA5, 0xFA)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
YELLOW = RGBColor(0xF1, 0xC4, 0x0F)
PURPLE = RGBColor(0x9B, 0x59, 0xB6)
TEAL = RGBColor(0x1A, 0xBC, 0x9C)
DARK_CARD = RGBColor(0x2D, 0x2D, 0x44)
DARKER_CARD = RGBColor(0x23, 0x23, 0x3A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # Smaller corner radius
    shape.adjustments[0] = 0.05
    return shape


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_arrow(slide, start_left, start_top, end_left, end_top, color=WHITE, width=Pt(2)):
    connector = slide.shapes.add_connector(
        1,  # straight connector
        start_left, start_top,
        end_left, end_top
    )
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def add_circle(slide, left, top, size, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape


def add_x_mark(slide, cx, cy, size=Inches(0.2), color=RED):
    """Draw an X mark using two lines."""
    half = size // 2
    # Line 1: top-left to bottom-right
    line1 = slide.shapes.add_connector(1, cx - half, cy - half, cx + half, cy + half)
    line1.line.color.rgb = color
    line1.line.width = Pt(3)
    # Line 2: top-right to bottom-left
    line2 = slide.shapes.add_connector(1, cx + half, cy - half, cx - half, cy + half)
    line2.line.color.rgb = color
    line2.line.width = Pt(3)


def add_check_mark(slide, cx, cy, color=GREEN):
    add_text(slide, cx - Inches(0.15), cy - Inches(0.2), Inches(0.4), Inches(0.4),
             "✓", font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 1: THE PROBLEM
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide1, DARK_BG)

# Title
add_text(slide1, Inches(0.5), Inches(0.25), Inches(12), Inches(0.6),
         "MHA PDS — Health Service Event Bundle Issue", font_size=28, color=WHITE, bold=True)
add_text(slide1, Inches(0.5), Inches(0.8), Inches(12), Inches(0.4),
         "Only 1 of 4 Health Service Events reaching Ontario Health", font_size=16, color=RED_LIGHT)

# Divider
add_rect(slide1, Inches(0.5), Inches(1.25), Inches(12.3), Pt(2), BLUE)

# --- LEFT SIDE: Cerner JSON Payload ---
add_text(slide1, Inches(0.5), Inches(1.45), Inches(5), Inches(0.4),
         "Cerner JSON Payload (4 Service Events)", font_size=16, color=BLUE_LIGHT, bold=True)

# SUBMIT_BUNDLE box
add_shape(slide1, Inches(0.5), Inches(1.95), Inches(5.5), Inches(1.6), DARK_CARD, BLUE, Pt(1))
add_text(slide1, Inches(0.7), Inches(2.0), Inches(5), Inches(0.35),
         "SUBMIT_BUNDLE (3 entries)", font_size=13, color=BLUE_LIGHT, bold=True)

# Bundle pills
bundle_y = Inches(2.4)
# EOC
add_shape(slide1, Inches(0.7), bundle_y, Inches(1.6), Inches(0.35), GREEN_DARK)
add_text(slide1, Inches(0.7), bundle_y, Inches(1.6), Inches(0.35),
         "SERVICE_REQUEST_EOC", font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
# SDOH
add_shape(slide1, Inches(2.45), bundle_y, Inches(1.3), Inches(0.35), PURPLE)
add_text(slide1, Inches(2.45), bundle_y, Inches(1.3), Inches(0.35),
         "CLIENT_SDOH", font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
# HSE - single
add_shape(slide1, Inches(3.9), bundle_y, Inches(1.9), Inches(0.35), RED)
add_text(slide1, Inches(3.9), bundle_y, Inches(1.9), Inches(0.35),
         "HEALTH_SERVICES ×1", font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Single UUID callout
add_text(slide1, Inches(0.7), Inches(2.9), Inches(5), Inches(0.5),
         "↑ 1 bundle = 1 UUID stamped on ALL 4 service rows", font_size=11, color=RED_LIGHT)

# SERVICES array
add_shape(slide1, Inches(0.5), Inches(3.65), Inches(5.5), Inches(3.3), DARK_CARD, ORANGE, Pt(1))
add_text(slide1, Inches(0.7), Inches(3.7), Inches(5), Inches(0.35),
         "SERVICES[] (4 service events)", font_size=13, color=ORANGE, bold=True)

svc_data = [
    ("SRV-3393360169", "2026-01-21", "In-person", "60 min"),
    ("SRV-3393361999", "2026-02-13", "Virtual", "30 min"),
    ("SRV-3393362203", "2026-02-13", "In-person", "30 min"),
    ("SRV-3393362296", "2026-02-13", "In-person", "60 min"),
]

for i, (eid, date, mode, mins) in enumerate(svc_data):
    sy = Inches(4.15) + Inches(i * 0.7)
    border = RED if i == 0 else MED_GRAY
    fill = DARKER_CARD if i > 0 else DARK_CARD
    add_shape(slide1, Inches(0.7), sy, Inches(5.1), Inches(0.6), fill, border, Pt(1.5))
    idx_color = RED_LIGHT if i == 0 else MED_GRAY
    add_text(slide1, Inches(0.85), sy + Inches(0.05), Inches(0.4), Inches(0.25),
             f"[{i}]", font_size=11, color=idx_color, bold=True)
    add_text(slide1, Inches(1.25), sy + Inches(0.05), Inches(2.2), Inches(0.25),
             eid, font_size=10, color=WHITE if i == 0 else LIGHT_GRAY)
    add_text(slide1, Inches(1.25), sy + Inches(0.3), Inches(4), Inches(0.25),
             f"{date}  •  {mode}  •  {mins}", font_size=9, color=LIGHT_GRAY if i == 0 else MED_GRAY)

    if i == 0:
        add_text(slide1, Inches(4.6), sy + Inches(0.15), Inches(1.2), Inches(0.3),
                 "← USED", font_size=11, color=RED_LIGHT, bold=True)
    else:
        add_text(slide1, Inches(4.6), sy + Inches(0.15), Inches(1.2), Inches(0.3),
                 "LOST", font_size=11, color=MED_GRAY, bold=True)

# --- MIDDLE: Arrow ---
add_text(slide1, Inches(6.15), Inches(3.8), Inches(1), Inches(0.5),
         "→", font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# --- RIGHT SIDE: Mirth Processing ---
add_text(slide1, Inches(7.2), Inches(1.45), Inches(5.5), Inches(0.4),
         "Mirth FHIR Transformation", font_size=16, color=BLUE_LIGHT, bold=True)

# Code issue box
add_shape(slide1, Inches(7.2), Inches(1.95), Inches(5.6), Inches(2.3), DARK_CARD, RED, Pt(1))
add_text(slide1, Inches(7.4), Inches(2.0), Inches(5), Inches(0.35),
         "Root Cause: Hardcoded SERVICES[0]", font_size=13, color=RED_LIGHT, bold=True)

# Code snippets
code_items = [
    ("createEncounterProfile()", "msg.SERVICES[0].HEALTH_SERVICE_EVENT"),
    ("createLocationProfile()", "msg.SERVICES[0].HSP_SITE"),
    ("createOrganizationProfile()", "msg.SERVICES[0].HSP_ORGANIZATION"),
]
for i, (func, code) in enumerate(code_items):
    cy = Inches(2.5) + Inches(i * 0.55)
    add_shape(slide1, Inches(7.4), cy, Inches(5.2), Inches(0.45), DARKER_CARD)
    add_text(slide1, Inches(7.55), cy + Inches(0.02), Inches(2.5), Inches(0.2),
             func, font_size=10, color=YELLOW, bold=True)
    add_text(slide1, Inches(7.55), cy + Inches(0.22), Inches(5), Inches(0.2),
             code, font_size=9, color=RED_LIGHT)

# Profile cache issue
add_shape(slide1, Inches(7.2), Inches(4.4), Inches(5.6), Inches(1.1), DARK_CARD, ORANGE, Pt(1))
add_text(slide1, Inches(7.4), Inches(4.45), Inches(5), Inches(0.35),
         "Profile Cache: Created once, reused for all", font_size=13, color=ORANGE, bold=True)
add_text(slide1, Inches(7.4), Inches(4.85), Inches(5.2), Inches(0.55),
         'getProfileCache("Encounter") → returns cached SERVICES[0] data\n'
         'Even in a loop, Encounter/Location/Org profiles never refresh',
         font_size=10, color=LIGHT_GRAY)

# Result box
add_shape(slide1, Inches(7.2), Inches(5.7), Inches(5.6), Inches(1.4), DARKER_CARD, RED, Pt(2))
add_text(slide1, Inches(7.4), Inches(5.75), Inches(5), Inches(0.35),
         "Result: Ontario Health Receives", font_size=14, color=WHITE, bold=True)

result_items = [
    ("1× SERVICE_REQUEST_EOC", GREEN, "✓"),
    ("1× CLIENT_SDOH", GREEN, "✓"),
    ("1× HEALTH_SERVICES (should be 4)", RED_LIGHT, "✗"),
]
for i, (label, color, icon) in enumerate(result_items):
    ry = Inches(6.15) + Inches(i * 0.3)
    add_text(slide1, Inches(7.6), ry, Inches(0.3), Inches(0.3),
             icon, font_size=14, color=color, bold=True)
    add_text(slide1, Inches(7.95), ry, Inches(4.5), Inches(0.3),
             label, font_size=12, color=color)


# ============================================================
# SLIDE 2: THE FIX
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, DARK_BG)

add_text(slide2, Inches(0.5), Inches(0.25), Inches(12), Inches(0.6),
         "Resolution — Per-Service HEALTH_SERVICES Bundles", font_size=28, color=WHITE, bold=True)
add_text(slide2, Inches(0.5), Inches(0.8), Inches(12), Inches(0.4),
         "Generate unique FHIR bundle per service event with independent tracking", font_size=16, color=GREEN)

add_rect(slide2, Inches(0.5), Inches(1.25), Inches(12.3), Pt(2), BLUE)

# --- LEFT: CCL Changes ---
add_text(slide2, Inches(0.5), Inches(1.45), Inches(5.5), Inches(0.4),
         "CCL Changes (Cerner Side)", font_size=16, color=BLUE_LIGHT, bold=True)

# Before box
add_shape(slide2, Inches(0.5), Inches(1.95), Inches(5.8), Inches(2.2), DARK_CARD, RED, Pt(1))
add_text(slide2, Inches(0.7), Inches(2.0), Inches(3), Inches(0.3),
         "BEFORE — sDetermine_Submit_Bundles", font_size=11, color=RED_LIGHT, bold=True)

add_shape(slide2, Inches(0.7), Inches(2.35), Inches(5.4), Inches(0.75), DARKER_CARD)
add_text(slide2, Inches(0.85), Inches(2.38), Inches(5.2), Inches(0.7),
         "SUBMIT_BUNDLE_CNT: 3\n"
         "  [1] SERVICE_REQUEST_EOC  →  UUID-A\n"
         "  [2] CLIENT_SDOH          →  UUID-B\n"
         "  [3] HEALTH_SERVICES ×1   →  UUID-C  (all 4 rows get UUID-C)",
         font_size=9, color=LIGHT_GRAY)

add_text(slide2, Inches(0.7), Inches(3.2), Inches(3), Inches(0.3),
         "AFTER — 1 bundle per service event", font_size=11, color=GREEN, bold=True)

add_shape(slide2, Inches(0.7), Inches(3.5), Inches(5.4), Inches(0.55), DARKER_CARD)
add_text(slide2, Inches(0.85), Inches(3.52), Inches(5.2), Inches(0.55),
         "SUBMIT_BUNDLE_CNT: 6\n"
         "  [1] SERVICE_REQUEST_EOC          →  UUID-A\n"
         "  [2] CLIENT_SDOH                  →  UUID-B",
         font_size=9, color=LIGHT_GRAY)

# HSE entries in green
hse_entries = [
    ("[3] HEALTH_SERVICES  SERVICE_IDX=0  →  UUID-C1",
     "SRV-3393360169"),
    ("[4] HEALTH_SERVICES  SERVICE_IDX=1  →  UUID-C2",
     "SRV-3393361999"),
    ("[5] HEALTH_SERVICES  SERVICE_IDX=2  →  UUID-C3",
     "SRV-3393362203"),
    ("[6] HEALTH_SERVICES  SERVICE_IDX=3  →  UUID-C4",
     "SRV-3393362296"),
]

for i, (entry, svc_id) in enumerate(hse_entries):
    ey = Inches(4.15) + Inches(i * 0.32)
    add_text(slide2, Inches(0.85), ey, Inches(5.2), Inches(0.3),
             f"  {entry}", font_size=9, color=GREEN)

# Unique UUID callout
add_text(slide2, Inches(0.7), Inches(5.5), Inches(5.5), Inches(0.5),
         "Each service row gets its own UUID for callback matching",
         font_size=11, color=GREEN, bold=True)

# Files changed
add_shape(slide2, Inches(0.5), Inches(5.95), Inches(5.8), Inches(1.2), DARK_CARD, BLUE, Pt(1))
add_text(slide2, Inches(0.7), Inches(6.0), Inches(5), Inches(0.3),
         "CCL Files Modified", font_size=11, color=BLUE_LIGHT, bold=True)

ccl_files = [
    "gbin_mha_pds_utils.prg → sDetermine_Submit_Bundles (N bundles, SERVICE_INDEX)",
    "gbin_mha_pds_data.prg  → per-service UUID assignment & status updates",
]
for i, f in enumerate(ccl_files):
    add_text(slide2, Inches(0.85), Inches(6.3) + Inches(i * 0.3), Inches(5.4), Inches(0.3),
             f"• {f}", font_size=9, color=LIGHT_GRAY)

# --- RIGHT: Mirth Changes ---
add_text(slide2, Inches(7.0), Inches(1.45), Inches(5.8), Inches(0.4),
         "Mirth Changes (FHIR Transformation)", font_size=16, color=BLUE_LIGHT, bold=True)

# SUBMIT_BUNDLE loop - no change needed
add_shape(slide2, Inches(7.0), Inches(1.95), Inches(5.8), Inches(1.1), DARK_CARD, GREEN, Pt(1))
add_text(slide2, Inches(7.2), Inches(2.0), Inches(5), Inches(0.3),
         "createBundleList() — No structural change needed", font_size=11, color=GREEN, bold=True)
add_text(slide2, Inches(7.2), Inches(2.35), Inches(5.4), Inches(0.6),
         "Already loops over SUBMIT_BUNDLE array.\n"
         "With 6 entries (instead of 3), it naturally creates 6 FHIR bundles.",
         font_size=10, color=LIGHT_GRAY)

# Profile functions fix
add_shape(slide2, Inches(7.0), Inches(3.2), Inches(5.8), Inches(2.2), DARK_CARD, ORANGE, Pt(1))
add_text(slide2, Inches(7.2), Inches(3.25), Inches(5), Inches(0.3),
         "Profile Functions — Use SERVICE_INDEX", font_size=11, color=ORANGE, bold=True)

fix_items = [
    ("createEncounterProfile()", "SERVICES[0]", "SERVICES[serviceIdx]"),
    ("createLocationProfile()", "SERVICES[0]", "SERVICES[serviceIdx]"),
    ("createOrganizationProfile()", "SERVICES[0]", "SERVICES[serviceIdx]"),
]

for i, (func, before, after) in enumerate(fix_items):
    fy = Inches(3.65) + Inches(i * 0.55)
    add_shape(slide2, Inches(7.2), fy, Inches(5.4), Inches(0.45), DARKER_CARD)
    add_text(slide2, Inches(7.35), fy + Inches(0.02), Inches(2.5), Inches(0.2),
             func, font_size=10, color=YELLOW, bold=True)
    add_text(slide2, Inches(7.35), fy + Inches(0.22), Inches(2.3), Inches(0.2),
             f"msg.{before}", font_size=9, color=RED_LIGHT)
    add_text(slide2, Inches(9.7), fy + Inches(0.22), Inches(0.3), Inches(0.2),
             "→", font_size=10, color=WHITE, bold=True)
    add_text(slide2, Inches(10.0), fy + Inches(0.22), Inches(2.5), Inches(0.2),
             f"msg.{after}", font_size=9, color=GREEN)

# Cache fix
add_shape(slide2, Inches(7.0), Inches(5.55), Inches(5.8), Inches(1.0), DARK_CARD, ORANGE, Pt(1))
add_text(slide2, Inches(7.2), Inches(5.6), Inches(5), Inches(0.3),
         "Profile Cache — Key per service index", font_size=11, color=ORANGE, bold=True)
add_text(slide2, Inches(7.2), Inches(5.95), Inches(5.4), Inches(0.5),
         'Cache key changes:  "Encounter" → "Encounter_0", "Encounter_1", etc.\n'
         'Ensures each service event gets its own Encounter/Location/Org profiles',
         font_size=10, color=LIGHT_GRAY)

# Result box
add_shape(slide2, Inches(7.0), Inches(6.7), Inches(5.8), Inches(0.55), DARKER_CARD, GREEN, Pt(2))
add_text(slide2, Inches(7.2), Inches(6.72), Inches(5.5), Inches(0.5),
         "Result: 1 EOC  +  1 SDOH  +  4 HEALTH_SERVICES  =  6 FHIR Bundles  ✓",
         font_size=13, color=GREEN, bold=True)


# ============================================================
# SLIDE 3: BEFORE vs AFTER FLOW
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, DARK_BG)

add_text(slide3, Inches(0.5), Inches(0.25), Inches(12), Inches(0.6),
         "Before vs After — End-to-End Data Flow", font_size=28, color=WHITE, bold=True)

add_rect(slide3, Inches(0.5), Inches(0.85), Inches(12.3), Pt(2), BLUE)

# === BEFORE (top half) ===
add_shape(slide3, Inches(0.3), Inches(1.1), Inches(1.2), Inches(0.4), RED, border_color=RED)
add_text(slide3, Inches(0.3), Inches(1.1), Inches(1.2), Inches(0.4),
         "BEFORE", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Cerner DB
add_shape(slide3, Inches(0.3), Inches(1.7), Inches(2.2), Inches(1.9), DARK_CARD, MED_GRAY, Pt(1))
add_text(slide3, Inches(0.3), Inches(1.7), Inches(2.2), Inches(0.35),
         "Cerner SERVICE Table", font_size=11, color=BLUE_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)

svc_labels = ["SRV-..0169", "SRV-..1999", "SRV-..2203", "SRV-..2296"]
for i, label in enumerate(svc_labels):
    sy = Inches(2.1) + Inches(i * 0.35)
    add_shape(slide3, Inches(0.45), sy, Inches(1.9), Inches(0.28), DARKER_CARD)
    add_text(slide3, Inches(0.45), sy, Inches(1.3), Inches(0.28),
             f" {label}", font_size=8, color=LIGHT_GRAY)
    # Same UUID marker
    add_text(slide3, Inches(1.65), sy, Inches(0.7), Inches(0.28),
             "UUID-C", font_size=7, color=RED_LIGHT, alignment=PP_ALIGN.RIGHT)

# Arrow
add_text(slide3, Inches(2.6), Inches(2.3), Inches(0.5), Inches(0.4),
         "→", font_size=24, color=MED_GRAY, bold=True)

# JSON payload
add_shape(slide3, Inches(3.1), Inches(1.7), Inches(3.0), Inches(1.9), DARK_CARD, MED_GRAY, Pt(1))
add_text(slide3, Inches(3.1), Inches(1.7), Inches(3.0), Inches(0.35),
         "JSON Payload", font_size=11, color=BLUE_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)

json_items = [
    ("SUBMIT_BUNDLE[3]", "3 bundles (1 HSE)", MED_GRAY),
    ("SERVICES[4]", "4 events in array", ORANGE),
]
for i, (key, val, col) in enumerate(json_items):
    jy = Inches(2.15) + Inches(i * 0.6)
    add_text(slide3, Inches(3.25), jy, Inches(2.7), Inches(0.25),
             key, font_size=9, color=col, bold=True)
    add_text(slide3, Inches(3.25), jy + Inches(0.22), Inches(2.7), Inches(0.25),
             val, font_size=8, color=LIGHT_GRAY)

# Arrow
add_text(slide3, Inches(6.2), Inches(2.3), Inches(0.5), Inches(0.4),
         "→", font_size=24, color=MED_GRAY, bold=True)

# Mirth
add_shape(slide3, Inches(6.7), Inches(1.7), Inches(2.8), Inches(1.9), DARK_CARD, MED_GRAY, Pt(1))
add_text(slide3, Inches(6.7), Inches(1.7), Inches(2.8), Inches(0.35),
         "Mirth FHIR Transform", font_size=11, color=BLUE_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide3, Inches(6.85), Inches(2.15), Inches(2.5), Inches(0.25),
         "Loops 3 SUBMIT_BUNDLEs", font_size=9, color=LIGHT_GRAY)
add_text(slide3, Inches(6.85), Inches(2.45), Inches(2.5), Inches(0.25),
         "Encounter → SERVICES[0] only", font_size=9, color=RED_LIGHT, bold=True)
add_text(slide3, Inches(6.85), Inches(2.75), Inches(2.5), Inches(0.25),
         "Cache locks to first profile", font_size=9, color=RED_LIGHT)

# Arrow
add_text(slide3, Inches(9.6), Inches(2.3), Inches(0.5), Inches(0.4),
         "→", font_size=24, color=MED_GRAY, bold=True)

# Ontario Health
add_shape(slide3, Inches(10.1), Inches(1.7), Inches(2.8), Inches(1.9), DARK_CARD, RED, Pt(2))
add_text(slide3, Inches(10.1), Inches(1.7), Inches(2.8), Inches(0.35),
         "Ontario Health", font_size=11, color=RED_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)

oh_before = [
    ("✓ 1× EOC bundle", GREEN),
    ("✓ 1× SDOH bundle", GREEN),
    ("✓ 1× HSE bundle", RED_LIGHT),
    ("✗ 3× HSE missing", RED),
]
for i, (label, color) in enumerate(oh_before):
    oy = Inches(2.15) + Inches(i * 0.33)
    add_text(slide3, Inches(10.25), oy, Inches(2.5), Inches(0.3),
             label, font_size=10, color=color)

# === AFTER (bottom half) ===
add_rect(slide3, Inches(0.3), Inches(3.85), Inches(12.7), Pt(1), MED_GRAY)

add_shape(slide3, Inches(0.3), Inches(4.1), Inches(1.2), Inches(0.4), GREEN_DARK, border_color=GREEN_DARK)
add_text(slide3, Inches(0.3), Inches(4.1), Inches(1.2), Inches(0.4),
         "AFTER", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Cerner DB
add_shape(slide3, Inches(0.3), Inches(4.7), Inches(2.2), Inches(2.3), DARK_CARD, GREEN, Pt(1))
add_text(slide3, Inches(0.3), Inches(4.7), Inches(2.2), Inches(0.35),
         "Cerner SERVICE Table", font_size=11, color=BLUE_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)

uuid_labels = ["UUID-C1", "UUID-C2", "UUID-C3", "UUID-C4"]
for i, (label, uuid) in enumerate(zip(svc_labels, uuid_labels)):
    sy = Inches(5.15) + Inches(i * 0.42)
    add_shape(slide3, Inches(0.45), sy, Inches(1.9), Inches(0.33), DARKER_CARD, GREEN, Pt(1))
    add_text(slide3, Inches(0.5), sy + Inches(0.02), Inches(1.2), Inches(0.28),
             f" {label}", font_size=8, color=WHITE)
    add_text(slide3, Inches(1.6), sy + Inches(0.02), Inches(0.7), Inches(0.28),
             uuid, font_size=7, color=GREEN, alignment=PP_ALIGN.RIGHT)

# Arrow
add_text(slide3, Inches(2.6), Inches(5.5), Inches(0.5), Inches(0.4),
         "→", font_size=24, color=GREEN, bold=True)

# JSON payload
add_shape(slide3, Inches(3.1), Inches(4.7), Inches(3.0), Inches(2.3), DARK_CARD, GREEN, Pt(1))
add_text(slide3, Inches(3.1), Inches(4.7), Inches(3.0), Inches(0.35),
         "JSON Payload", font_size=11, color=BLUE_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide3, Inches(3.25), Inches(5.15), Inches(2.7), Inches(0.25),
         "SUBMIT_BUNDLE[6]", font_size=9, color=GREEN, bold=True)
add_text(slide3, Inches(3.25), Inches(5.37), Inches(2.7), Inches(0.2),
         "6 bundles (4 HSE, each w/ index)", font_size=8, color=LIGHT_GRAY)

bundle_list = [
    ("EOC", GREEN_DARK, "seq 1"),
    ("SDOH", PURPLE, "seq 2"),
    ("HSE idx=0", TEAL, "seq 3"),
    ("HSE idx=1", TEAL, "seq 4"),
    ("HSE idx=2", TEAL, "seq 5"),
    ("HSE idx=3", TEAL, "seq 6"),
]
for i, (label, col, seq) in enumerate(bundle_list):
    by = Inches(5.65) + Inches(i * 0.22)
    add_text(slide3, Inches(3.35), by, Inches(2.7), Inches(0.2),
             f"{seq}: {label}", font_size=8, color=col if i >= 2 else LIGHT_GRAY)

# Arrow
add_text(slide3, Inches(6.2), Inches(5.5), Inches(0.5), Inches(0.4),
         "→", font_size=24, color=GREEN, bold=True)

# Mirth
add_shape(slide3, Inches(6.7), Inches(4.7), Inches(2.8), Inches(2.3), DARK_CARD, GREEN, Pt(1))
add_text(slide3, Inches(6.7), Inches(4.7), Inches(2.8), Inches(0.35),
         "Mirth FHIR Transform", font_size=11, color=BLUE_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide3, Inches(6.85), Inches(5.15), Inches(2.5), Inches(0.25),
         "Loops 6 SUBMIT_BUNDLEs", font_size=9, color=LIGHT_GRAY)
add_text(slide3, Inches(6.85), Inches(5.45), Inches(2.5), Inches(0.25),
         "SERVICE_INDEX → correct", font_size=9, color=GREEN, bold=True)
add_text(slide3, Inches(6.85), Inches(5.65), Inches(2.5), Inches(0.25),
         "  SERVICES[idx]", font_size=9, color=GREEN)
add_text(slide3, Inches(6.85), Inches(5.95), Inches(2.5), Inches(0.25),
         "Cache keyed per index", font_size=9, color=GREEN)
add_text(slide3, Inches(6.85), Inches(6.25), Inches(2.5), Inches(0.25),
         "Unique profiles per event", font_size=9, color=LIGHT_GRAY)

# Arrow
add_text(slide3, Inches(9.6), Inches(5.5), Inches(0.5), Inches(0.4),
         "→", font_size=24, color=GREEN, bold=True)

# Ontario Health
add_shape(slide3, Inches(10.1), Inches(4.7), Inches(2.8), Inches(2.3), DARK_CARD, GREEN, Pt(2))
add_text(slide3, Inches(10.1), Inches(4.7), Inches(2.8), Inches(0.35),
         "Ontario Health", font_size=11, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

oh_after = [
    "✓ 1× EOC bundle",
    "✓ 1× SDOH bundle",
    "✓ 1× HSE (SRV-..0169)",
    "✓ 1× HSE (SRV-..1999)",
    "✓ 1× HSE (SRV-..2203)",
    "✓ 1× HSE (SRV-..2296)",
]
for i, label in enumerate(oh_after):
    oy = Inches(5.15) + Inches(i * 0.28)
    add_text(slide3, Inches(10.25), oy, Inches(2.5), Inches(0.25),
             label, font_size=10, color=GREEN)

add_text(slide3, Inches(10.25), Inches(6.9), Inches(2.5), Inches(0.3),
         "6 bundles total  ✓", font_size=12, color=GREEN, bold=True)

# Save
output_path = "/Users/chadcummings/Github/chadcumm/gbin-mha-pds/public/temp/MHA_PDS_HSE_Bundle_Issue.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
